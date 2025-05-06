from flask import Flask, request, jsonify
import requests
import io
import pdfplumber
from pptx import Presentation

app = Flask(__name__)

@app.route("/analyzeFile", methods=["POST"])
def analyze_file():
    data = request.json
    item_id = data["itemId"]
    file_type = data["fileType"]

    access_token = "HIER_DEIN_ACCESS_TOKEN"

    url = f"https://graph.microsoft.com/v1.0/me/drive/items/{item_id}/content"
    headers = {"Authorization": f"Bearer {access_token}"}
    r = requests.get(url, headers=headers)
    file = io.BytesIO(r.content)

    if file_type == "pdf":
        result = extract_text_from_pdf(file)
    elif file_type == "pptx":
        result = extract_text_from_pptx(file)
    else:
        return jsonify({"error": "Unsupported file type"}), 400

    return jsonify(result)

def extract_text_from_pdf(file):
    text = ""
    with pdfplumber.open(file) as pdf:
        for i, page in enumerate(pdf.pages):
            page_text = page.extract_text()
            text += f"\n--- Seite {i+1} ---\n{page_text}"
    return {"text": text}

def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = ""
    for i, slide in enumerate(prs.slides):
        text += f"\n--- Folie {i+1} ---\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return {"text": text}

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=10000, debug=True)
